"""
文档加载器封装

统一接口加载不同格式的文档。
支持格式：PDF, Word, Markdown, HTML, TXT, CSV, Excel, PPT, 图片（OCR）
"""

import os
from typing import List, Optional
from pathlib import Path

from app.core.config import settings


class Document:
    """加载后的文档对象"""

    def __init__(
        self,
        content: str,
        metadata: Optional[dict] = None,
    ):
        self.content = content
        self.metadata = metadata or {}

    def __repr__(self):
        title = self.metadata.get("title", "unknown")
        size = len(self.content) if self.content else 0
        return f"<Document title='{title}' size={size}>"


class BaseLoader:
    """文档加载器基类"""

    def load(self, file_path: str) -> Document:
        """加载文档
        Args:
            file_path: 文件路径（本地或 MinIO 下载后的临时路径）
        Returns:
            Document 对象
        """
        raise NotImplementedError


class PDFLoader(BaseLoader):
    """PDF 文档加载器

    使用 pypdf 提取文本，支持多页 PDF。
    如需更好的格式保留，可升级为 PyMuPDF (fitz)。
    """

    def load(self, file_path: str) -> Document:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ImportError("pypdf not installed. Run: pip install pypdf")

        try:
            reader = PdfReader(file_path)
        except Exception as e:
            raise ValueError(f"无法解析 PDF 文件（文件可能已损坏或不是有效的 PDF 格式）: {e}")

        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text)

        content = "\n\n".join(pages)
        if not content.strip():
            raise ValueError("PDF 文件内容为空或无法提取文本")

        return Document(
            content=content,
            metadata={
                "title": Path(file_path).stem,
                "page_count": len(reader.pages),
                "file_type": "pdf",
            },
        )


class DocxLoader(BaseLoader):
    """Word 文档加载器"""

    def load(self, file_path: str) -> Document:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # 也提取表格中的文本
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)

        content = "\n\n".join(paragraphs)
        return Document(
            content=content,
            metadata={
                "title": Path(file_path).stem,
                "paragraph_count": len(paragraphs),
                "file_type": "docx",
            },
        )


class MarkdownLoader(BaseLoader):
    """Markdown 文档加载器

    Markdown 内容直接读取，保留原始格式。
    标题结构和格式标记对后续 heading-aware 分块有信号作用。
    """

    def load(self, file_path: str) -> Document:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()

        # 尝试从第一个 # 标题提取文档标题
        title = Path(file_path).stem
        for line in content.split("\n"):
            if line.startswith("# "):
                title = line[2:].strip()
                break

        return Document(
            content=content,
            metadata={
                "title": title,
                "file_type": "markdown",
            },
        )


class HTMLLoader(BaseLoader):
    """HTML 文档加载器

    使用 BeautifulSoup 提取文本内容。
    移除 script / style 标签，保留段落结构。
    """

    def load(self, file_path: str) -> Document:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")

        with open(file_path, "r", encoding="utf-8") as f:
            html = f.read()

        soup = BeautifulSoup(html, "html.parser")

        # 移除无用标签
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        # 提取文本
        text = soup.get_text(separator="\n")

        # 清理多余空行
        import re
        text = re.sub(r'\n{3,}', '\n\n', text)

        # 提取标题
        title = Path(file_path).stem
        if soup.title:
            title = soup.title.get_text(strip=True)

        return Document(
            content=text.strip(),
            metadata={
                "title": title,
                "file_type": "html",
            },
        )


class TxtLoader(BaseLoader):
    """纯文本文档加载器（多编码尝试）"""

    def load(self, file_path: str) -> Document:
        # 按优先级尝试不同编码
        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            raise ValueError(f"Cannot decode file: {file_path}")

        return Document(
            content=content,
            metadata={
                "title": Path(file_path).stem,
                "file_type": "txt",
            },
        )


class CSVLoader(BaseLoader):
    """CSV 文件加载器"""

    def load(self, file_path: str) -> Document:
        import csv
        with open(file_path, "r", encoding="utf-8-sig") as f:
            reader = csv.reader(f)
            rows = list(reader)

        if not rows:
            return Document(content="", metadata={"title": Path(file_path).stem, "file_type": "csv"})

        header = rows[0]
        lines = []
        if header:
            lines.append(" | ".join(header))
            lines.append("-" * len(lines[0]))
        for row in rows[1:]:
            lines.append(" | ".join(row))

        return Document(
            content="\n".join(lines) if lines else "",
            metadata={
                "title": Path(file_path).stem,
                "row_count": len(rows),
                "file_type": "csv",
            },
        )


class ExcelLoader(BaseLoader):
    """Excel (.xlsx) 加载器"""

    def load(self, file_path: str) -> Document:
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise ImportError("openpyxl not installed. Run: pip install openpyxl")

        wb = load_workbook(file_path, read_only=True, data_only=True)
        all_sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            all_sheets.append(f"## Sheet: {sheet_name}")
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    all_sheets.append(row_text)
            all_sheets.append("")
        wb.close()

        return Document(
            content="\n".join(all_sheets),
            metadata={
                "title": Path(file_path).stem,
                "sheet_count": len(wb.sheetnames),
                "file_type": "xlsx",
            },
        )


class PPTLoader(BaseLoader):
    """PowerPoint (.pptx) 加载器"""

    def load(self, file_path: str) -> Document:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            slide_lines = [f"## Slide {i + 1}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            slide_lines.append(row_text)
            slides_text.append("\n".join(slide_lines))

        return Document(
            content="\n\n".join(slides_text),
            metadata={
                "title": Path(file_path).stem,
                "slide_count": len(prs.slides),
                "file_type": "pptx",
            },
        )


# ─── 加载器工厂 ───

class ImageLoader(BaseLoader):
    """Image OCR loader via DashScope OmniParser / Qwen-VL-OCR."""

    def load(self, file_path: str) -> Document:
        import base64
        try:
            import dashscope
            from http import HTTPStatus
        except ImportError:
            raise ImportError("dashscope not installed. Run: pip install dashscope")

        with open(file_path, "rb") as f:
            image_data = base64.b64encode(f.read()).decode("utf-8")

        messages = [{
            "role": "user",
            "content": [
                {"image": f"data:image/png;base64,{image_data}"},
                {"text": "请提取图片中的所有文字内容，保持原有格式，不要添加额外说明。"},
            ],
        }]

        resp = dashscope.MultiModalConversation.call(
            model="qwen-vl-max",
            messages=messages,
            api_key=settings.EMBEDDING_API_KEY,
        )
        if resp.status_code != HTTPStatus.OK:
            raise RuntimeError(f"OCR API error: {resp.code} {resp.message}")

        # Parse OCR response — format varies by model version
        output = resp.output
        if output and output.choices:
            content = output.choices[0].message.content
            if isinstance(content, str):
                text = content
            elif isinstance(content, list):
                texts = []
                for part in content:
                    if isinstance(part, dict):
                        texts.append(part.get("text", ""))
                    elif hasattr(part, "text"):
                        texts.append(part.text)
                text = "\n".join(t for t in texts if t)
            else:
                text = str(content)
        else:
            text = ""

        if not text.strip():
            raise RuntimeError("OCR returned empty text — image may have no text content")

        return Document(
            content=text,
            metadata={
                "title": Path(file_path).stem,
                "file_type": "image",
                "ocr_model": "qwen-vl-ocr",
            },
        )


LOADER_MAP = {
    "pdf": PDFLoader,
    "docx": DocxLoader,
    "doc": DocxLoader,
    "md": MarkdownLoader,
    "markdown": MarkdownLoader,
    "html": HTMLLoader,
    "htm": HTMLLoader,
    "txt": TxtLoader,
    "text": TxtLoader,
    "csv": CSVLoader,
    "xlsx": ExcelLoader,
    "xls": ExcelLoader,
    "pptx": PPTLoader,
    "ppt": PPTLoader,
    "png": ImageLoader,
    "jpg": ImageLoader,
    "jpeg": ImageLoader,
}


def load_document(file_path: str, file_type: Optional[str] = None) -> Document:
    """
    根据文件类型自动选择加载器

    Args:
        file_path: 文件路径
        file_type: 文件类型（如不提供则从扩展名推断）

    Returns:
        Document 对象

    Raises:
        ValueError: 不支持的文件类型
    """
    if file_type is None:
        file_type = Path(file_path).suffix.lstrip(".").lower()

    loader_class = LOADER_MAP.get(file_type)
    if loader_class is None:
        raise ValueError(
            f"Unsupported file type: {file_type}. "
            f"Supported types: {list(LOADER_MAP.keys())}"
        )

    loader = loader_class()
    return loader.load(file_path)
