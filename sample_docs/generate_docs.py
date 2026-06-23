"""Generate sample documents in PDF, TXT, MD, DOCX, PPTX formats."""
import os

BASE = os.path.dirname(os.path.abspath(__file__))

# ============================================================
# 1. TXT
# ============================================================
txt_content = (
    "====================================================================\n"
    "           KnowFlow 企业知识库平台 v2.0 正式发布公告\n"
    "====================================================================\n\n"
    "发布日期：2026 年 6 月 15 日\n"
    "发布部门：技术研发部\n"
    "文档编号：KF-REL-2026-001\n\n"
    "---\n\n"
    "一、版本概述\n\n"
    "KnowFlow v2.0 是面向企业级场景的 RAG 知识库问答平台，基于 LangGraph + Milvus\n"
    "构建。本次大版本升级引入了动态 RRF 混合检索、自适应三层分块、记忆压缩等核心能力。\n\n"
    "二、核心特性\n\n"
    "1. 动态 RRF 混合检索\n"
    "   - 基于查询意图自动调节向量/关键词权重\n"
    "   - 向量检索依托 Milvus 2.5 集群，支持十亿级数据秒级检索\n"
    "   - 关键词检索基于 BM25 算法，精准匹配专业术语\n\n"
    "2. 自适应三层分块\n"
    "   - Structural Chunking：识别 Markdown 标题、代码块边界\n"
    "   - Semantic Chunking：基于 embedding 相似度切分\n"
    "   - Fixed Chunking：兜底策略，确保每块不超过 512 tokens\n"
    "   - 支持 Parent-Child 模式，检索用小块、返回用大块\n\n"
    "3. 记忆压缩 (Memory Consolidation)\n"
    "   - LLM 驱动的对话历史摘要\n"
    "   - Token 感知的自动触发（超过 32K tokens 时压缩）\n"
    "   - 压缩率最高 90%，保留关键信息不丢失\n\n"
    "三、性能指标\n\n"
    "+------------------+-----------+------------+\n"
    "| 指标              | v1.0      | v2.0       |\n"
    "+------------------+-----------+------------+\n"
    "| 检索延迟 (P50)    | 850ms     | 320ms      |\n"
    "| 检索延迟 (P99)    | 2.1s      | 0.8s       |\n"
    "| RRF 融合准确率    | 78%       | 91%        |\n"
    "| 单文档处理时间     | 12s       | 4s         |\n"
    "| 支持最大 Chunk 数 | 100 万    | 5000 万    |\n"
    "+------------------+-----------+------------+\n\n"
    "四、升级指南\n\n"
    "从 v1.x 升级到 v2.0 的步骤：\n"
    "1. 备份现有 PostgreSQL 数据库\n"
    "2. 拉取最新镜像: docker compose pull\n"
    "3. 执行数据库迁移: alembic upgrade head\n"
    "4. 重启服务: docker compose up -d\n"
    "5. 验证服务健康状态: curl http://localhost:8000/health\n\n"
    "五、已知问题\n"
    "- 超大 PDF（>500 页）建议拆分为多个文档上传\n"
    "- 图片 OCR 功能将在 v2.1 版本支持\n"
    "- Windows 环境下 Milvus 需要 WSL2 后端\n\n"
    "六、联系方式\n"
    "技术支持：support@knowflow.io\n"
    "文档地址：https://docs.knowflow.io\n"
)

with open(os.path.join(BASE, "KnowFlow_v2.0_发布公告.txt"), "w", encoding="utf-8") as f:
    f.write(txt_content)
print("1/5  TXT done")

# ============================================================
# 2. MD
# ============================================================
md_content = """\
# Kubernetes 集群运维手册

> **适用范围**: KnowFlow 生产环境
> **版本**: v1.3
> **维护者**: 基础架构组

---

## 1. 集群架构

### 节点列表

| 节点名称 | IP | 角色 | 状态 |
|---------|-----|------|------|
| kf-master-01 | 10.0.1.10 | Control Plane | Running |
| kf-worker-01 | 10.0.1.21 | Backend | Running |
| kf-worker-02 | 10.0.1.22 | Backend | Running |
| kf-worker-03 | 10.0.1.23 | Milvus | Running |
| kf-worker-04 | 10.0.1.24 | PostgreSQL | Running |

---

## 2. 日常巡检

### 2.1 检查集群健康状态

```bash
kubectl get nodes
kubectl get pods -n knowflow
kubectl get svc -n knowflow
kubectl get pvc -n knowflow
```

### 2.2 检查资源使用

```bash
kubectl top nodes
kubectl top pods -n knowflow --sort-by=cpu
```

---

## 3. 常见故障处理

### 3.1 Pod 重启循环 (CrashLoopBackOff)

**排查步骤**:

1. 查看 Pod 事件: `kubectl describe pod <pod-name> -n knowflow`
2. 查看 Error 日志: `kubectl logs <pod-name> -n knowflow | grep -i error | tail -20`
3. 常见原因:
   - PostgreSQL / Milvus 连接失败 → 检查 Service 和 Endpoint
   - OOM → 调整 resources.limits.memory
   - 配置文件错误 → 检查 ConfigMap

### 3.2 Milvus 检索延迟升高

1. 检查 Index 状态: `curl http://localhost:9091/healthz`
2. 检查磁盘 IO: `iostat -x 1 5`
3. 临时措施：触发索引重建

### 3.3 PostgreSQL 主从延迟

```sql
SELECT * FROM pg_stat_replication;
SELECT application_name, state, sync_state,
  pg_wal_lsn_diff(pg_current_wal_lsn(), sent_lsn) AS sent_lag,
  pg_wal_lsn_diff(pg_current_wal_lsn(), write_lsn) AS write_lag
FROM pg_stat_replication;
```

---

## 4. 备份与恢复

### 数据库备份

```bash
kubectl exec -it deployment/knowflow-postgres -n knowflow -- \\
  pg_dump -U knowflow -Fc knowflow > knowflow_$(date +%Y%m%d).dump
```

### 恢复流程

```bash
# 1. 停止写入
kubectl scale deployment knowflow-backend -n knowflow --replicas=0
# 2. 恢复数据库
kubectl exec -it deployment/knowflow-postgres -n knowflow -- \\
  pg_restore -U knowflow -d knowflow --clean knowflow_20260623.dump
# 3. 重建索引
# 4. 恢复服务
kubectl scale deployment knowflow-backend -n knowflow --replicas=3
```

---

## 5. 扩容指南

```bash
# Backend 水平扩容 (HPA)
kubectl autoscale deployment knowflow-backend -n knowflow \\
  --cpu-percent=70 --min=3 --max=20

# Milvus DataNode 扩容
kubectl scale deployment knowflow-milvus-datanode -n knowflow --replicas=5
```

> **最后更新**: 2026-06-23
> **下次审查**: 2026-09-23
"""

with open(os.path.join(BASE, "K8s_运维手册.md"), "w", encoding="utf-8") as f:
    f.write(md_content)
print("2/5  MD done")

# ============================================================
# 3. DOCX
# ============================================================
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

doc = Document()
doc.add_heading("新员工入职指南", level=0).alignment = WD_ALIGN_PARAGRAPH.CENTER
doc.add_paragraph("欢迎加入 KnowFlow 团队！本指南将帮助你在入职第一周快速融入。")

doc.add_heading("1. 公司介绍", level=1)
doc.add_paragraph(
    "KnowFlow 是一家专注于企业知识管理的人工智能公司。"
    "我们的核心产品是企业级 RAG 知识库问答平台，帮助组织高效管理和检索内部知识。"
    "公司成立于 2024 年，目前已服务超过 200 家企业客户。"
)

doc.add_heading("2. 入职准备", level=1)
doc.add_heading("2.1 硬件设备", level=2)

table = doc.add_table(rows=5, cols=3, style="Light Grid Accent 1")
for i, h in enumerate(["设备", "型号", "领取地点"]):
    table.rows[0].cells[i].text = h
for i, row_data in enumerate([
    ["笔记本电脑", "MacBook Pro M4 32GB", "IT 部 3 楼"],
    ["显示器", "Dell 27\" 4K", "IT 部 3 楼"],
    ["耳机", "AirPods Pro", "行政部 2 楼"],
    ["工卡", "-", "前台"],
]):
    for j, cell_text in enumerate(row_data):
        table.rows[i + 1].cells[j].text = cell_text

doc.add_heading("2.2 账号开通", level=2)
for a in [
    "企业邮箱：hr@knowflow.io",
    "GitHub：加入 KnowFlow 组织",
    "飞书：扫描工卡背面二维码",
    "VPN：安装 WireGuard，配置文件由 IT 提供",
]:
    doc.add_paragraph(a, style="List Bullet")

doc.add_heading("3. 开发环境搭建", level=1)
doc.add_paragraph(
    "克隆代码仓库: git clone git@github.com:LLIEsmallwhite/KnowFlow.git\n"
    "安装依赖: uv sync\n"
    "配置环境: cp .env.example .env\n"
    "启动服务: docker compose up -d\n"
    "初始化数据库: cd backend && alembic upgrade head\n"
    "开发服务器: uvicorn main:app --reload --port 8000"
)

doc.add_heading("4. 第一天任务清单", level=1)
for t in [
    "完成设备领取和账号开通",
    "克隆代码仓库，完成环境搭建",
    "阅读 README.md 和 CLAUDE.md",
    "成功运行 docker compose up -d",
    "上传一份测试文档到知识库",
    "提交第一个 PR（修改 README 中的错别字）",
    "参加下午 3 点的新人欢迎会",
]:
    doc.add_paragraph(t, style="List Bullet")

doc.add_heading("5. 团队通讯录", level=1)
ct = doc.add_table(rows=5, cols=4, style="Light Grid Accent 1")
for i, h in enumerate(["姓名", "职位", "飞书", "GitHub"]):
    ct.rows[0].cells[i].text = h
for i, row_data in enumerate([
    ["张伟", "技术负责人", "@zhangwei", "@zhangwei-dev"],
    ["李娜", "产品经理", "@lina", "@lina-pm"],
    ["王强", "后端工程师", "@wangqiang", "@wangqiang"],
    ["陈晓", "前端工程师", "@chenxiao", "@chenxiao"],
]):
    for j, ct_text in enumerate(row_data):
        ct.rows[i + 1].cells[j].text = ct_text

doc.add_paragraph("")
doc.add_paragraph("如有任何问题，随时在飞书群 #new-hires 提问！", style="Intense Quote")
doc.save(os.path.join(BASE, "新员工入职指南.docx"))
print("3/5  DOCX done")

# ============================================================
# 4. PPTX
# ============================================================
from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides_data = [
    ("2026 Q2 技术季度总结", [
        "汇报人：技术研发部 | 日期：2026-06-23",
        "KnowFlow — 企业级 RAG 知识库平台",
    ]),
    ("里程碑达成", [
        "KnowFlow v2.0 发布 — 动态 RRF + 自适应分块",
        "混合检索延迟降低 62%（850ms -> 320ms P50）",
        "RRF 融合准确率提升至 91%（+13pp）",
        "通过 ISO 27001 信息安全认证",
        "客户数突破 200 家，日活用户 > 5000",
    ]),
    ("核心技术创新", [
        "动态 RRF 权重计算 — 基于查询语义自动调配向量/关键词权重",
        "  代码/ID 类查询 -> 关键词权重 up",
        "  语义/概念类查询 -> 向量权重 up",
        "三层自适应分块: Structural + Semantic + Fixed",
        "LLM 记忆压缩 — 32K 阈值触发，压缩率最高 90%",
    ]),
    ("性能优化成果", [
        "Milvus 索引: IVF_FLAT -> IVF_SQ8，内存节省 75%",
        "连接池 + 查询缓存，并发能力提升 3x",
        "Reranker: API 调用 -> 本地模型，延迟降低 90%",
        "Celery 异步任务池，支持 20 个并发解析",
        "大文件分片上传，支持 500MB+ 文档",
    ]),
    ("Q3 规划", [
        "KnowFlow v2.1 — 目标 9 月底发布",
        "  多模态检索: 图片 OCR + 表格理解",
        "  联网搜索: Bing/Google Search API",
        "  Agent 模式: 多步推理 + 工具调用",
        "  权限体系: RBAC + 知识库级隔离",
        "性能目标: P50 < 200ms, P99 < 500ms",
    ]),
    ("团队成长", [
        "Q2 加入 8 位新成员（4 后端 + 2 前端 + 1 算法 + 1 QA）",
        "技术分享 12 场（RRF 算法、Milvus 调优、LangGraph 实战...）",
        "开源贡献: 向 LangChain/Milvus 提交 5 个 PR",
        "团队飞书群活跃度 up 45%",
    ]),
]

for title, bullets in slides_data:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1]
    for b in bullets:
        body.text_frame.add_paragraph().text = b

prs.save(os.path.join(BASE, "Q2技术季度总结.pptx"))
print("4/5  PPTX done")

# ============================================================
# 5. PDF
# ============================================================
try:
    from fpdf import FPDF
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Courier", size=10)
    lines = [
        "KnowFlow API Interface Documentation",
        "Version: v2.0",
        "Date: 2026-06-23",
        "",
        "=" * 60,
        "",
        "1. Knowledge Base API",
        "  POST   /api/v1/knowledge-bases          Create KB",
        "  GET    /api/v1/knowledge-bases          List KBs",
        "  DELETE /api/v1/knowledge-bases/{id}     Delete KB",
        "  POST   /api/v1/knowledge-bases/{id}/documents  Upload doc",
        "",
        "2. Chat API",
        "  POST   /api/v1/chat                     Sync Q&A",
        "  POST   /api/v1/chat/stream              SSE streaming Q&A",
        "",
        "3. Agent API",
        "  POST   /api/v1/agent/chat               Agent Q&A",
        "  POST   /api/v1/agent/chat/stream        Agent streaming",
        "",
        "4. Auth API",
        "  POST   /api/v1/auth/register            Register",
        "  POST   /api/v1/auth/login               Login",
        "  GET    /api/v1/auth/me                  Current user",
        "",
        "5. Conversations API",
        "  GET    /api/v1/conversations            List sessions",
        "  POST   /api/v1/conversations            Create session",
        "  DELETE /api/v1/conversations/{id}       Delete session",
        "",
        "Error Codes: 400 401 404 409 500",
        "=" * 60,
        "Total endpoints: 22 | KnowFlow v2.0",
    ]
    for line in lines:
        pdf.cell(200, 6, txt=line[:100], ln=True)
    pdf.output(os.path.join(BASE, "KnowFlow_API文档.pdf"))
    print("5/5  PDF done (fpdf2)")
except ImportError:
    # Fallback: write as text file with .pdf extension for testing
    pdf_text = (
        "KnowFlow API Interface Documentation\n"
        "Version: v2.0 | Date: 2026-06-23\n"
        "=" * 60 + "\n\n"
        "1. Knowledge Base API\n"
        "  POST   /api/v1/knowledge-bases          Create KB\n"
        "  GET    /api/v1/knowledge-bases          List KBs\n"
        "  DELETE /api/v1/knowledge-bases/{id}     Delete KB\n"
        "  POST   /api/v1/knowledge-bases/{id}/documents  Upload doc\n\n"
        "2. Chat API\n"
        "  POST   /api/v1/chat                     Sync Q&A\n"
        "  POST   /api/v1/chat/stream              SSE streaming\n\n"
        "3. Agent API\n"
        "  POST   /api/v1/agent/chat               Agent Q&A\n\n"
        "4. Auth API\n"
        "  POST   /api/v1/auth/register / login\n\n"
        "5. Conversations API\n"
        "  GET/POST/DELETE /api/v1/conversations\n\n"
        "Error Codes: 400 401 404 409 500\n"
        "Total endpoints: 22"
    )
    with open(os.path.join(BASE, "KnowFlow_API文档.pdf"), "w", encoding="utf-8") as f:
        f.write(pdf_text)
    print("5/5  PDF done (plain text - fpdf2 not available)")

print("\nAll 5 sample documents generated in sample_docs/")
